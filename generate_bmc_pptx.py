import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_bmc_presentation(output_pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    C_NAVY_DARK  = RGBColor(15, 23, 42)    # #0F172A
    C_PRIMARY    = RGBColor(0, 79, 159)   # #004F9F (SUS Blue)
    C_TEAL       = RGBColor(0, 165, 153)  # #00A599 (Accent Teal)
    C_SLATE_BG   = RGBColor(248, 250, 252)# #F8FAFC
    C_WHITE      = RGBColor(255, 255, 255)
    C_TEXT_DARK  = RGBColor(30, 41, 59)   # #1E293B
    C_TEXT_MUTED = RGBColor(100, 116, 139)# #64748B
    C_BORDER     = RGBColor(226, 232, 240)# #E2E8F0

    logo_path = os.path.join(os.getcwd(), "logo_digitrilha.png")

    # SLIDE 1
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_NAVY_DARK
    bg1.line.color.rgb = C_NAVY_DARK

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.12), Inches(5.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TEAL
    bar.line.color.rgb = C_TEAL

    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(1.2), Inches(1.0), width=Inches(3.8))

    tb_title = slide1.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(10.5), Inches(4.2))
    tf1 = tb_title.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p_badge = tf1.paragraphs[0]
    p_badge.text = "HACKATHON SAÚDE DIGITAL UNIFESP 2026 • DESAFIO 2"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = C_TEAL
    p_badge.font.name = "Arial"
    p_badge.space_after = Pt(10)

    p_main = tf1.add_paragraph()
    p_main.text = "Business Model Canvas"
    p_main.font.size = Pt(42)
    p_main.font.bold = True
    p_main.font.color.rgb = C_WHITE
    p_main.font.name = "Arial"

    p_sub = tf1.add_paragraph()
    p_sub.text = "Digitrilha: Capacitação Permanente em Tecnologia Assistiva (Órteses e Próteses)"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(203, 213, 225)
    p_sub.font.name = "Arial"
    p_sub.space_before = Pt(6)
    p_sub.space_after = Pt(20)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Modelagem de negócios orientada ao impacto no SUS, viabilidade técnica, sustentabilidade financeira e segurança de dados (LGPD)."
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = C_TEXT_MUTED
    p_desc.font.name = "Arial"

    # SLIDE 2
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_SLATE_BG
    bg2.line.color.rgb = C_SLATE_BG

    header_box = slide2.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.533), Inches(0.65))
    tf_h = header_box.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    p_h.text = "Business Model Canvas — Digitrilha"
    p_h.font.size = Pt(20)
    p_h.font.bold = True
    p_h.font.color.rgb = C_PRIMARY
    p_h.font.name = "Arial"
    
    p_hs = tf_h.add_paragraph()
    p_hs.text = "Plataforma de Microlearning Adaptativo e Suporte Clínico por IA para Atenção Primária do SUS"
    p_hs.font.size = Pt(11)
    p_hs.font.color.rgb = C_TEXT_MUTED
    p_hs.font.name = "Arial"

    left_m = Inches(0.4)
    top_m = Inches(0.9)
    col_w = Inches(2.426)
    gap = Inches(0.09)
    
    top_h = Inches(4.45)
    half_h = Inches(2.18)

    bot_top = Inches(5.44)
    bot_h = Inches(1.85)

    def draw_canvas_block(slide, title, items, left, top, width, height, accent_color=C_PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(1)

        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.35))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent_color
        strip.line.color.rgb = accent_color
        
        tb_hdr = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.07), width - Inches(0.2), Inches(0.3))
        tf_hdr = tb_hdr.text_frame
        tf_hdr.word_wrap = True
        tf_hdr.margin_left = tf_hdr.margin_top = tf_hdr.margin_right = tf_hdr.margin_bottom = 0
        p_hdr = tf_hdr.paragraphs[0]
        p_hdr.text = title.upper()
        p_hdr.font.size = Pt(9.5)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = C_WHITE
        p_hdr.font.name = "Arial"
        p_hdr.alignment = PP_ALIGN.LEFT

        tb_body = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.45), width - Inches(0.2), height - Inches(0.5))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True
        tf_body.margin_left = tf_body.margin_top = tf_body.margin_right = tf_body.margin_bottom = 0

        for idx, item in enumerate(items):
            p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(8.5)
            p.font.color.rgb = C_TEXT_DARK
            p.font.name = "Arial"
            p.space_after = Pt(2.5)

    draw_canvas_block(slide2, "1. Parcerias-Chave", ["Ministério da Saúde (SGTES): fomento & homologação no AVASUS.", "UNIFESP & Universidades: curadoria científica em Órteses e Próteses.", "CERs (Centros Esp. Reabilitação): apoio a simulações e feedback de fluxo.", "CONASEMS / COSEMS: articulação junto a Secretarias Municipais.", "Conselhos (CREFITO, CFM, COFEN): chancela e horas de EPS."], left_m, top_m, col_w, top_h, C_PRIMARY)
    draw_canvas_block(slide2, "2. Atividades-Chave", ["Curadoria pedagógica de microlearning em O&P (3-5 min).", "Treinamento RAG/IA da assistente Eva com portarias do SUS.", "Desenvolvimento de simulações RPG 2D e WebXR VR.", "Integração de APIs (Gov.br, e-SUS APS, AVASUS).", "Monitoramento de engajamento e métricas de impacto."], left_m + col_w + gap, top_m, col_w, half_h, C_PRIMARY)
    draw_canvas_block(slide2, "3. Recursos-Chave", ["Base de conhecimento em O&P e normativas MS.", "Motor IA Eva (LLM RAG + guardrails LGPD).", "Plataforma Web/PWA leve (uso offline 3G/4G).", "Equipe multidisciplinar (Fisiatria, IA, Edu, UX)."], left_m + col_w + gap, top_m + half_h + gap, col_w, half_h, C_PRIMARY)
    draw_canvas_block(slide2, "4. Proposta de Valor", ["Profissionais APS: capacitação prática diária (3-5 min) + assistente IA 24/7 para apoio à decisão clínica em O&P.", "Gestores SUS: aumento da resolutividade na UBS, redução de filas no CER e otimização do gasto público.", "Pacientes (PCD/Mobilidade): atendimento ágil, prescrição assertiva, menor abandono e prevenção de escaras.", "Rede SUS: interoperabilidade (e-SUS/Gov.br), conformidade LGPD e alta escalabilidade nacional."], left_m + (col_w + gap)*2, top_m, col_w, top_h, C_TEAL)
    draw_canvas_block(slide2, "5. Relacionamento", ["Eva Tutor IA: tirada de dúvidas 24/7 e feedback clínico imediato.", "Gamificação: XP, badges (ex: Guardião da Postura) e micro-certificados.", "Suporte contínuo & canais de feedback pedagógico."], left_m + (col_w + gap)*3, top_m, col_w, half_h, C_PRIMARY)
    draw_canvas_block(slide2, "6. Canais", ["Plataformas Oficiais: AVASUS e UNA-SUS.", "e-SUS APS: módulo DSS (apoio à decisão no prontuário).", "PWA Digitrilha: Web/Mobile (compatível offline).", "Eventos CONASEMS/COSEMS & redes InterPET."], left_m + (col_w + gap)*3, top_m + half_h + gap, col_w, half_h, C_PRIMARY)
    draw_canvas_block(slide2, "7. Segmentos de Clientes", ["Clientes Compradores (B2G): Ministério da Saúde, Secretarias Estaduais e Municipais de Saúde (SES/SMS).", "Usuários Finais (Profissionais APS): Médicos ESF, Fisioterapeutas, Terapeutas Ocupacionais, Enfermeiros e ACS.", "Beneficiários Finais (Pacientes): Pessoas com deficiência (PCD) ou mobilidade reduzida atendidas pelo SUS."], left_m + (col_w + gap)*4, top_m, col_w, top_h, C_PRIMARY)
    
    bot_w1 = col_w * 3 + gap * 2
    draw_canvas_block(slide2, "8. Estrutura de Custos", ["Infraestrutura Cloud & LLM API: hospedagem serverless, consumo de IA e CDN.", "Desenvolvimento de Software & Jogos: manutenção do PWA, engine RPG/VR e APIs de integração.", "Curadoria & Conteúdo Clínico: remuneração de especialistas de saúde e instrucionais.", "Governança, Suporte & LGPD: atendimento B2G, auditorias de segurança e treinamento."], left_m, bot_top, bot_w1, bot_h, RGBColor(180, 83, 9))

    bot_w2 = col_w * 2 + gap
    draw_canvas_block(slide2, "9. Fontes de Receita / Financiamento", ["Licenciamento B2G / EPS: contratos anuais com municípios/estados (Fundo Nac. de Saúde).", "Editais de Inovação & Fomento: FINEP, CNPq, DECIT/MS e programas de transformação digital.", "Customização Regional: adaptação de fluxos regulatórios específicos por estado/município.", "Parcerias ESG / Impacto: fomento institucional para ampliação de trilhas educacionais."], left_m + bot_w1 + gap, bot_top, bot_w2, bot_h, RGBColor(16, 185, 129))

    # SLIDE 3
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = C_WHITE
    bg3.line.color.rgb = C_WHITE

    tb3_title = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf3 = tb3_title.text_frame
    tf3.word_wrap = True
    p3_t = tf3.paragraphs[0]
    p3_t.text = "Detalhamento Estratégico do Digitrilha"
    p3_t.font.size = Pt(26)
    p3_t.font.bold = True
    p3_t.font.color.rgb = C_PRIMARY

    p3_sub = tf3.add_paragraph()
    p3_sub.text = "Como o modelo de negócios responde diretamente aos critérios de avaliação do Hackathon Unifesp 2026"
    p3_sub.font.size = Pt(14)
    p3_sub.font.color.rgb = C_TEXT_MUTED

    pillars = [
        ("Impacto Clínico & Social (30%)", C_PRIMARY, ["Empoderamento da Atenção Primária para triagem e prescrição rápida.", "Redução da taxa de abandono de cadeiras de rodas e órteses.", "Prevenção de complicações graves (lesões por pressão e deformidades)."]),
        ("Viabilidade Técnica (30%)", C_TEAL, ["Arquitetura modular em microserviços integrável via LTI ao AVASUS.", "Suporte à decisão clínica (DSS) conectável ao e-SUS APS.", "Autenticação segura e unificada via Gov.br + CNES."]),
        ("Segurança & LGPD (20%)", RGBColor(124, 58, 237), ["Uso exclusivo de dados sintéticos nos simuladores clínicos.", "Criptografia ponta a ponta (TLS) no tráfego de dados.", "Termo de Consentimento Livre e Esclarecido (TCLE) digital."]),
        ("Escalabilidade Nacional (20%)", RGBColor(16, 185, 129), ["PWA otimizado para baixas conexões (3G/4G) e funcionamento offline.", "Flexibilidade para adaptação dos fluxos regulatórios regionais.", "Capacitação contínua sem necessidade de afastar o profissional da UBS."])
    ]

    p_w = Inches(2.7)
    p_h = Inches(4.5)
    p_gap = Inches(0.3)

    for i, (p_title, p_color, p_bullets) in enumerate(pillars):
        px = Inches(0.8) + i * (p_w + p_gap)
        py = Inches(1.8)

        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, p_w, p_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_SLATE_BG
        card.line.color.rgb = C_BORDER

        top_strip = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, p_w, Inches(0.8))
        top_strip.fill.solid()
        top_strip.fill.fore_color.rgb = p_color
        top_strip.line.color.rgb = p_color

        tb_pt = slide3.shapes.add_textbox(px + Inches(0.15), py + Inches(0.1), p_w - Inches(0.3), Inches(0.6))
        tf_pt = tb_pt.text_frame
        tf_pt.word_wrap = True
        ppt = tf_pt.paragraphs[0]
        ppt.text = p_title
        ppt.font.size = Pt(13)
        ppt.font.bold = True
        ppt.font.color.rgb = C_WHITE

        tb_pb = slide3.shapes.add_textbox(px + Inches(0.15), py + Inches(0.9), p_w - Inches(0.3), p_h - Inches(1.0))
        tf_pb = tb_pb.text_frame
        tf_pb.word_wrap = True

        for b_idx, bullet in enumerate(p_bullets):
            p = tf_pb.paragraphs[0] if b_idx == 0 else tf_pb.add_paragraph()
            p.text = f"✔ {bullet}"
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(10)

    try:
        prs.save(output_pptx_path)
        print(f"Apresentação PPTX gerada com sucesso em: {output_pptx_path}")
    except PermissionError:
        alt_path = output_pptx_path.replace(".pptx", "_final.pptx")
        prs.save(alt_path)
        print(f"Arquivo principal bloqueado pelo usuário. Apresentação PPTX salva em: {alt_path}")

if __name__ == "__main__":
    output_path = os.path.join(os.getcwd(), "business_model_canvas_digitrilha.pptx")
    create_bmc_presentation(output_path)
